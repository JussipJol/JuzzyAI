from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree

# Color palette
BG_DARK = RGBColor(0x0D, 0x11, 0x17)       # very dark navy
BG_CARD = RGBColor(0x13, 0x1A, 0x26)       # dark card
ACCENT = RGBColor(0x00, 0xD4, 0xFF)         # cyan accent
ACCENT2 = RGBColor(0x7C, 0x3A, 0xED)        # purple
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0xA0, 0xAE, 0xBF)
GREEN = RGBColor(0x00, 0xE5, 0x96)
YELLOW = RGBColor(0xFF, 0xD6, 0x00)

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]  # blank

def add_slide():
    slide = prs.slides.add_slide(blank_layout)
    # background rectangle
    bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_DARK
    bg.line.fill.background()
    return slide

def txb(slide, text, left, top, width, height,
        font_size=24, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
        italic=False, font_name="Segoe UI"):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return tb

def rect(slide, left, top, width, height, fill_color, line_color=None, radius=False):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape

def accent_bar(slide, top=0.55):
    bar = rect(slide, 0, top, 13.33, 0.04, ACCENT)
    return bar

def add_multiline_txb(slide, lines, left, top, width, height,
                       font_size=18, color=WHITE, bold_first=False, font_name="Segoe UI"):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.bold = (bold_first and i == 0)
        run.font.name = font_name
    return tb

# ─────────────────────────────────────────────
# SLIDE 1 — Title
# ─────────────────────────────────────────────
slide = add_slide()

# gradient-ish left panel
panel = rect(slide, 0, 0, 6.5, 7.5, RGBColor(0x13, 0x1A, 0x26))

# cyan vertical stripe
rect(slide, 0, 0, 0.08, 7.5, ACCENT)

# big logo-text
txb(slide, "JuzzyAI", 0.3, 1.5, 6, 1.5, font_size=72, bold=True, color=ACCENT, font_name="Segoe UI Black")
txb(slide, "AI-помощник разработчика\nпрямо в терминале", 0.3, 3.3, 6, 1.5, font_size=28, color=WHITE)
txb(slide, "v2.0.0 Prototype", 0.3, 5.2, 3, 0.5, font_size=16, color=GRAY, italic=True)

# right side — tagline cards
rect(slide, 7, 1.2, 5.8, 1.1, RGBColor(0x00, 0x44, 0x66), line_color=ACCENT)
txb(slide, "💬  Чат с AI", 7.2, 1.35, 5.4, 0.8, font_size=22, color=WHITE)

rect(slide, 7, 2.5, 5.8, 1.1, RGBColor(0x1A, 0x0A, 0x44), line_color=ACCENT2)
txb(slide, "🔍  Анализ кода", 7.2, 2.65, 5.4, 0.8, font_size=22, color=WHITE)

rect(slide, 7, 3.8, 5.8, 1.1, RGBColor(0x00, 0x3D, 0x2B), line_color=GREEN)
txb(slide, "⚡  Генерация кода", 7.2, 3.95, 5.4, 0.8, font_size=22, color=WHITE)

rect(slide, 7, 5.1, 5.8, 1.1, RGBColor(0x3D, 0x2B, 0x00), line_color=YELLOW)
txb(slide, "🔧  Рефакторинг", 7.2, 5.25, 5.4, 0.8, font_size=22, color=WHITE)

txb(slide, "Бесплатно • Офлайн • Кроссплатформенно", 7, 6.5, 5.8, 0.6,
    font_size=14, color=GRAY, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────
# SLIDE 2 — Проблема и решение
# ─────────────────────────────────────────────
slide = add_slide()
accent_bar(slide)

txb(slide, "Проблема и наше решение", 0.5, 0.1, 12, 0.55,
    font_size=32, bold=True, color=WHITE)

# Problem side
rect(slide, 0.4, 0.8, 5.8, 5.9, RGBColor(0x1A, 0x08, 0x08), line_color=RGBColor(0xFF, 0x44, 0x44))
txb(slide, "❌  Как раньше", 0.7, 0.9, 5, 0.6, font_size=22, bold=True, color=RGBColor(0xFF, 0x66, 0x66))

problems = [
    "• Постоянно переключаться между IDE,\n   браузером, ChatGPT",
    "• Копипастить код вручную",
    "• Платные подписки: $20–$50/мес",
    "• Интернет обязателен",
    "• Код уходит на сторонние серверы",
]
y = 1.6
for pr in problems:
    txb(slide, pr, 0.7, y, 5.3, 0.7, font_size=17, color=RGBColor(0xFF, 0xAA, 0xAA))
    y += 0.82

# Solution side
rect(slide, 6.8, 0.8, 6.1, 5.9, RGBColor(0x04, 0x1A, 0x0E), line_color=GREEN)
txb(slide, "✅  С JuzzyAI", 7.1, 0.9, 5.5, 0.6, font_size=22, bold=True, color=GREEN)

solutions = [
    "• Всё в одном терминале — без переключений",
    "• AI сам читает и редактирует файлы",
    "• Полностью бесплатно (5 провайдеров)",
    "• Работает офлайн через Ollama",
    "• Код остаётся у вас локально",
]
y = 1.6
for sol in solutions:
    txb(slide, sol, 7.1, y, 5.6, 0.7, font_size=17, color=RGBColor(0xAA, 0xFF, 0xCC))
    y += 0.82

# ─────────────────────────────────────────────
# SLIDE 3 — Функционал (демо)
# ─────────────────────────────────────────────
slide = add_slide()
accent_bar(slide)
txb(slide, "Ключевые функции", 0.5, 0.1, 12, 0.55, font_size=32, bold=True, color=WHITE)

features = [
    ("💬", "Умный чат", "Задавайте вопросы по коду, получайте объяснения.\nИстория сессий сохраняется.", ACCENT, RGBColor(0x00, 0x22, 0x44)),
    ("🔍", "Анализ кода", "Находит баги, уязвимости, плохие практики\nв вашем файле за секунды.", ACCENT2, RGBColor(0x18, 0x08, 0x38)),
    ("⚡", "Генерация", "Опишите задачу — получите готовый код.\nAI создаёт файлы напрямую.", GREEN, RGBColor(0x02, 0x1A, 0x0E)),
    ("🔧", "Рефакторинг", "Чистит и улучшает ваш код.\nПоказывает diff до и после.", YELLOW, RGBColor(0x1F, 0x18, 0x00)),
]

cols = [(0.3, 3.1), (3.6, 3.1), (6.9, 3.1), (10.2, 2.8)]
y_start = 0.95

for i, (icon, title, desc, accent_c, bg_c) in enumerate(features):
    x = cols[i][0]
    w = cols[i][1]
    card = rect(slide, x, y_start, w, 5.8, bg_c, line_color=accent_c)
    txb(slide, icon, x + 0.15, y_start + 0.2, w - 0.2, 0.8, font_size=40)
    txb(slide, title, x + 0.15, y_start + 1.1, w - 0.2, 0.6,
        font_size=20, bold=True, color=accent_c)
    txb(slide, desc, x + 0.15, y_start + 1.85, w - 0.2, 1.5, font_size=16, color=WHITE)

# bottom bar — terminal hint
rect(slide, 0.3, 6.1, 12.7, 1.1, RGBColor(0x08, 0x08, 0x08), line_color=RGBColor(0x33, 0x33, 0x33))
txb(slide, "$ juzzyai chat     $ juzzyai analyze --file app.py     $ juzzyai generate     $ juzzyai refactor --file utils.py",
    0.5, 6.2, 12.3, 0.8, font_size=16, color=GREEN, font_name="Consolas")

# ─────────────────────────────────────────────
# SLIDE 4 — Провайдеры / Офлайн режим
# ─────────────────────────────────────────────
slide = add_slide()
accent_bar(slide)
txb(slide, "5 AI-провайдеров — всё бесплатно", 0.5, 0.1, 12, 0.55,
    font_size=32, bold=True, color=WHITE)

providers = [
    ("Ollama", "Локально\n(офлайн)", GREEN, "🖥️"),
    ("Groq", "Облако\nБесплатно", ACCENT, "☁️"),
    ("OpenRouter", "28+ моделей\nБесплатно", ACCENT2, "🌐"),
    ("Gemini", "Google AI\nБесплатный тир", RGBColor(0xFF, 0x88, 0x00), "🔶"),
    ("HuggingFace", "Open-source\nмодели", RGBColor(0xFF, 0xD6, 0x00), "🤗"),
]

xs = [0.3, 2.92, 5.54, 8.16, 10.78]
for i, (name, desc, color, icon) in enumerate(providers):
    x = xs[i]
    rect(slide, x, 0.9, 2.4, 3.2, RGBColor(0x13, 0x1A, 0x26), line_color=color)
    txb(slide, icon, x + 0.7, 1.0, 1.2, 0.8, font_size=32, align=PP_ALIGN.CENTER)
    txb(slide, name, x + 0.1, 1.9, 2.2, 0.5, font_size=18, bold=True, color=color, align=PP_ALIGN.CENTER)
    txb(slide, desc, x + 0.1, 2.5, 2.2, 1.2, font_size=15, color=GRAY, align=PP_ALIGN.CENTER)

# Big highlight: offline
rect(slide, 0.3, 4.4, 12.7, 2.6, RGBColor(0x02, 0x1A, 0x0E), line_color=GREEN)
txb(slide, "🌐  Работает полностью офлайн", 0.6, 4.55, 8, 0.6,
    font_size=26, bold=True, color=GREEN)
txb(slide, "Через Ollama вы запускаете AI-модель прямо на своём компьютере.\nНикаких утечек кода, никаких подписок, никакого интернета.",
    0.6, 5.2, 9, 1.3, font_size=18, color=WHITE)
txb(slide, "🔒 Ваш код\nникуда\nне уходит", 10.5, 4.5, 2.5, 2.2,
    font_size=20, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────
# SLIDE 5 — Конкуренты и отличия
# ─────────────────────────────────────────────
slide = add_slide()
accent_bar(slide)
txb(slide, "Аналоги и наши преимущества", 0.5, 0.1, 12, 0.55,
    font_size=32, bold=True, color=WHITE)

# Table header
headers = ["Функция", "GitHub Copilot", "ChatGPT", "Codeium", "JuzzyAI ✨"]
header_colors = [GRAY, GRAY, GRAY, GRAY, ACCENT]
col_w = [2.8, 2.1, 2.1, 2.1, 2.6]
col_x = [0.3]
for w in col_w[:-1]:
    col_x.append(col_x[-1] + w + 0.1)

y_header = 0.9
for i, (h, cx, cw) in enumerate(zip(headers, col_x, col_w)):
    bg = ACCENT if i == 4 else RGBColor(0x1A, 0x22, 0x33)
    rect(slide, cx, y_header, cw, 0.5, bg)
    tc = BG_DARK if i == 4 else WHITE
    txb(slide, h, cx + 0.05, y_header + 0.05, cw - 0.1, 0.45,
        font_size=15, bold=True, color=tc, align=PP_ALIGN.CENTER)

rows = [
    ["Цена",            "$10–19/мес",  "$20/мес",   "Частично",    "Бесплатно"],
    ["Офлайн режим",    "❌",           "❌",         "❌",           "✅"],
    ["Анализ файлов",   "Частично",    "❌",         "❌",           "✅ Полный"],
    ["Генерация файлов","Частично",    "❌",         "❌",           "✅ AI пишет"],
    ["Конфиденц-ть",    "⚠️ Облако",   "⚠️ Облако",  "⚠️ Облако",   "✅ Локально"],
    ["Плагины",         "❌",           "❌",         "❌",           "✅ /review и др"],
]

row_colors = [RGBColor(0x0D, 0x11, 0x17), RGBColor(0x10, 0x16, 0x1E)]
for ri, row in enumerate(rows):
    y = y_header + 0.55 + ri * 0.7
    for ci, (cell, cx, cw) in enumerate(zip(row, col_x, col_w)):
        bg = row_colors[ri % 2] if ci < 4 else RGBColor(0x00, 0x22, 0x11)
        rect(slide, cx, y, cw, 0.65, bg)
        tc = GREEN if (ci == 4 and ("✅" in cell or "Бесплатно" in cell or "Локально" in cell)) else (RGBColor(0xFF, 0x88, 0x88) if "❌" in cell else WHITE)
        txb(slide, cell, cx + 0.05, y + 0.08, cw - 0.1, 0.55,
            font_size=14, color=tc, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────
# SLIDE 6 — Монетизация и внедрение
# ─────────────────────────────────────────────
slide = add_slide()
accent_bar(slide)
txb(slide, "Бизнес-модель и внедрение", 0.5, 0.1, 12, 0.55,
    font_size=32, bold=True, color=WHITE)

# Left: monetization
rect(slide, 0.3, 0.8, 5.9, 6.1, RGBColor(0x13, 0x1A, 0x26), line_color=ACCENT2)
txb(slide, "💰 Монетизация", 0.55, 0.9, 5.4, 0.5, font_size=22, bold=True, color=ACCENT2)

money_items = [
    ("Freemium лицензии", "Базовая версия бесплатна.\nPro-лицензия: расширенные функции,\nприоритетная поддержка."),
    ("Enterprise", "Коробочная версия для компаний:\nустановка на сервер, контроль доступа,\nинтеграция с CI/CD."),
    ("Маркетплейс плагинов", "Сторонние разработчики создают\nплатные плагины для JuzzyAI."),
]
y = 1.55
for title, desc in money_items:
    txb(slide, "▸ " + title, 0.55, y, 5.3, 0.4, font_size=17, bold=True, color=ACCENT2)
    txb(slide, desc, 0.55, y + 0.42, 5.3, 0.9, font_size=15, color=GRAY)
    y += 1.45

# Right: deployment
rect(slide, 6.7, 0.8, 6.2, 6.1, RGBColor(0x13, 0x1A, 0x26), line_color=ACCENT)
txb(slide, "🚀 Внедрение", 6.95, 0.9, 5.7, 0.5, font_size=22, bold=True, color=ACCENT)

deploy_items = [
    ("Этап 1 — GitHub Releases", "Публикация .exe-установщика\nдля Windows. Уже реализовано."),
    ("Этап 2 — Колледжи и университеты", "Партнёрство с учебными\nзаведениями, пилот-программы."),
    ("Этап 3 — B2B продажи", "IT-компании, аутсорс-команды,\nстартапы. Коробочная версия."),
]
y = 1.55
for title, desc in deploy_items:
    txb(slide, "▸ " + title, 6.95, y, 5.7, 0.4, font_size=17, bold=True, color=ACCENT)
    txb(slide, desc, 6.95, y + 0.42, 5.7, 0.9, font_size=15, color=GRAY)
    y += 1.45

# Budget
rect(slide, 0.3, 6.0, 12.6, 0.9, RGBColor(0x1A, 0x14, 0x00), line_color=YELLOW)
txb(slide, "💵  Бюджет на запуск коробочной версии: ~150 000 – 300 000 тг  (серверы, домен, маркетинг, юр.оформление)",
    0.5, 6.1, 12.2, 0.7, font_size=16, color=YELLOW)

# ─────────────────────────────────────────────
# SLIDE 7 — Ответы на вопросы
# ─────────────────────────────────────────────
slide = add_slide()
accent_bar(slide)
txb(slide, "Ответы на вопросы жюри", 0.5, 0.1, 12, 0.55,
    font_size=32, bold=True, color=WHITE)

qa = [
    ("Кто что делал?",
     "Жуссуп — backend (AI-интеграция, провайдеры, архитектура).\nПартнёр — frontend UX, тестирование, документация."),
    ("Почему эти инструменты?",
     "Python — стандарт AI/ML экосистемы (LangChain, Ollama, Groq SDK).\nRich — лучший рендеринг Markdown в терминале. PyInstaller — .exe без зависимостей."),
    ("Кто консультировал по бизнес-процессам?",
     "Интервью с 12 разработчиками из IT-компаний Казахстана.\nОбратная связь по pain-points: переключение контекста, утечки кода, стоимость."),
    ("Перспективы развития",
     "GUI-версия, VS Code Extension, веб-дашборд для команд,\nмаркетплейс плагинов, поддержка казахского языка в промптах."),
]

col1 = [(0, 2), (2, 4)]  # rows for left col
col2 = [(1, 3), (3, 5)]

positions = [
    (0.3, 0.85, 6.2),
    (0.3, 3.25, 6.2),
    (6.8, 0.85, 6.2),
    (6.8, 3.25, 6.2),
]
box_h = 2.15

for i, (q, a) in enumerate(qa):
    x, y, w = positions[i]
    color = [ACCENT, ACCENT2, GREEN, YELLOW][i]
    rect(slide, x, y, w, box_h, RGBColor(0x13, 0x1A, 0x26), line_color=color)
    txb(slide, q, x + 0.15, y + 0.08, w - 0.2, 0.45,
        font_size=17, bold=True, color=color)
    txb(slide, a, x + 0.15, y + 0.58, w - 0.2, 1.45,
        font_size=14, color=WHITE)

# ─────────────────────────────────────────────
# SLIDE 8 — Итог / CTA
# ─────────────────────────────────────────────
slide = add_slide()

# full-width gradient panel
rect(slide, 0, 0, 13.33, 7.5, RGBColor(0x06, 0x0D, 0x18))
rect(slide, 0, 0, 0.12, 7.5, ACCENT)

txb(slide, "JuzzyAI", 0.4, 0.7, 10, 1.2,
    font_size=80, bold=True, color=ACCENT, font_name="Segoe UI Black")
txb(slide, "Ваш AI-напарник в терминале.", 0.4, 2.1, 10, 0.8,
    font_size=34, color=WHITE)

bullets = [
    "✅  Бесплатно — 5 AI-провайдеров",
    "✅  Офлайн — код остаётся у вас",
    "✅  Всё в одном: чат, анализ, генерация, рефакторинг",
    "✅  Плагины — расширяйте сами",
]
y = 3.1
for b in bullets:
    txb(slide, b, 0.4, y, 9, 0.55, font_size=22, color=WHITE)
    y += 0.65

# CTA box
rect(slide, 0.4, 6.0, 8, 1.1, ACCENT, line_color=ACCENT)
txb(slide, "github.com/Zhussup/juzzyai  •  Скачать бесплатно", 0.5, 6.1, 7.8, 0.9,
    font_size=20, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)

# right side decoration
txb(slide, "🚀", 10.5, 2.5, 2.5, 2.5, font_size=120, align=PP_ALIGN.CENTER)

out_path = r"c:\Users\User\Documents\JuzzyAI\JuzzyAI_Presentation.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
